from __future__ import annotations

import argparse
import json
import textwrap
import zipfile
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Iterable

import fitz
from PIL import Image, ImageColor, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer


SLIDE_WIDTH = 1080
SLIDE_HEIGHT = 1350
INSTA_RATIO = (10, 12.5)


def now_iso() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat().replace("+00:00", "Z")


def find_font(candidates: Iterable[str], size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    for candidate in candidates:
        path = Path(candidate)
        if path.exists():
            return ImageFont.truetype(str(path), size=size)
    return ImageFont.load_default()


TITLE_FONT = find_font(
    [
        "C:/Windows/Fonts/georgiab.ttf",
        "C:/Windows/Fonts/arialbd.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSerif-Bold.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    ],
    72,
)
BODY_FONT = find_font(
    [
        "C:/Windows/Fonts/arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ],
    34,
)
SMALL_FONT = find_font(
    [
        "C:/Windows/Fonts/arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ],
    24,
)


@dataclass(frozen=True)
class CarouselTheme:
    slug: str
    kicker: str
    title: str
    subtitle: str
    insight_title: str
    bullets: list[str]
    cta_title: str
    cta_lines: list[str]
    background: str
    accent: str
    text: str
    muted: str


THEMES: list[CarouselTheme] = [
    CarouselTheme(
        slug="positioning",
        kicker="Brand strategy",
        title="Sharpen your positioning in 3 moves",
        subtitle="A clean carousel set for creators who need clarity before they post.",
        insight_title="Use these prompts",
        bullets=["Who is this for?", "What outcome do they want?", "Why choose you first?"],
        cta_title="Turn this into content",
        cta_lines=["Slide 1: promise", "Slide 2: proof", "Slide 3: next step"],
        background="#F5F0EA",
        accent="#7F6A58",
        text="#1F1D1A",
        muted="#7B756F",
    ),
    CarouselTheme(
        slug="lead-magnet",
        kicker="Lead generation",
        title="Lead magnets that feel simple, not pushy",
        subtitle="Editorial layouts for soft-selling digital products and services.",
        insight_title="What to include",
        bullets=["Fast win", "Tangible outcome", "Clear next action"],
        cta_title="Template flow",
        cta_lines=["Problem", "Framework", "Offer"],
        background="#F8F4EF",
        accent="#B89A86",
        text="#22201E",
        muted="#8A8078",
    ),
    CarouselTheme(
        slug="content-engine",
        kicker="Content systems",
        title="Build a repeatable weekly content engine",
        subtitle="Neutral slides for creators who want more consistency without loud design.",
        insight_title="Weekly rhythm",
        bullets=["Teach one idea", "Show one proof point", "Invite one action"],
        cta_title="Batch faster",
        cta_lines=["Write once", "Remix twice", "Schedule ahead"],
        background="#F3EFE8",
        accent="#4E5661",
        text="#191C21",
        muted="#6E737B",
    ),
    CarouselTheme(
        slug="offer-stack",
        kicker="Offer design",
        title="Package a simple offer stack that sells",
        subtitle="Use polished carousels to explain the value behind your service or product.",
        insight_title="Core ingredients",
        bullets=["Outcome", "Process", "Boundaries"],
        cta_title="Make it skimmable",
        cta_lines=["Pain point", "Solution", "Price logic"],
        background="#FAF7F2",
        accent="#A38E78",
        text="#201E1B",
        muted="#857D75",
    ),
    CarouselTheme(
        slug="case-study",
        kicker="Proof",
        title="Turn small wins into clean case-study posts",
        subtitle="A calm visual system for before-and-after stories and client outcomes.",
        insight_title="Case-study structure",
        bullets=["Context", "Action", "Result"],
        cta_title="End with momentum",
        cta_lines=["Keep it specific", "Use one metric", "Invite the next step"],
        background="#F1EEE8",
        accent="#5D5A54",
        text="#1F1C18",
        muted="#7C766F",
    ),
    CarouselTheme(
        slug="newsletter",
        kicker="Audience growth",
        title="Create carousels that feed your newsletter",
        subtitle="Made for consultants and creators who need quiet, premium-looking promotion.",
        insight_title="Conversion notes",
        bullets=["Lead with a lesson", "Offer a deeper resource", "Repeat the CTA once"],
        cta_title="Keep the path obvious",
        cta_lines=["Read", "Save", "Subscribe"],
        background="#F9F6F1",
        accent="#998873",
        text="#1E1C19",
        muted="#7E776F",
    ),
    CarouselTheme(
        slug="service-educator",
        kicker="Service marketing",
        title="Teach the client before you pitch the service",
        subtitle="Use these slides to educate prospects and build trust before the offer.",
        insight_title="Teaching angles",
        bullets=["Common mistake", "Better framing", "Easy first fix"],
        cta_title="Subtle pitch pattern",
        cta_lines=["Educate", "Reframe", "Invite"],
        background="#F4F1EC",
        accent="#6D7580",
        text="#1D2127",
        muted="#727A83",
    ),
    CarouselTheme(
        slug="launch-week",
        kicker="Launch planning",
        title="Design a calmer launch week",
        subtitle="Minimal templates for pre-launch, launch-day, and recap sequences.",
        insight_title="Launch stack",
        bullets=["Build curiosity", "Name the offer", "Show urgency cleanly"],
        cta_title="Three-post arc",
        cta_lines=["Preview", "Launch", "Recap"],
        background="#FAF4EE",
        accent="#B68068",
        text="#241E1B",
        muted="#8F7A70",
    ),
    CarouselTheme(
        slug="productized-service",
        kicker="Productization",
        title="Make a productized service feel premium",
        subtitle="Explain what is included, what is not, and why the structure helps the buyer.",
        insight_title="What to clarify",
        bullets=["Scope", "Timeline", "Expected result"],
        cta_title="Reduce buyer friction",
        cta_lines=["Answer objections", "Show the flow", "Point to the next step"],
        background="#F7F3EE",
        accent="#62554A",
        text="#211D19",
        muted="#7C726A",
    ),
    CarouselTheme(
        slug="monthly-plan",
        kicker="Planning",
        title="Plan a month of content in one sitting",
        subtitle="Template sets for creators who want a flexible system, not a rigid calendar.",
        insight_title="Plan around these buckets",
        bullets=["Teach", "Show proof", "Invite action"],
        cta_title="Use the pack weekly",
        cta_lines=["Duplicate", "Edit copy", "Publish"],
        background="#F2EFEA",
        accent="#8E8D89",
        text="#1E1E1D",
        muted="#757573",
    ),
]


def wrap_lines(text: str, width: int) -> list[str]:
    return textwrap.wrap(text, width=width) or [text]


def to_rgb(hex_color: str) -> tuple[int, int, int]:
    return ImageColor.getrgb(hex_color)


def add_text_block(draw: ImageDraw.ImageDraw, *, text: str, xy: tuple[int, int], font, fill: str, max_width: int, line_gap: int) -> int:
    x, y = xy
    current_y = y
    for paragraph in text.split("\n"):
      lines = wrap_lines(paragraph, max(12, max_width // max(font.size, 10)))
      for line in lines:
        draw.text((x, current_y), line, font=font, fill=fill)
        bbox = draw.textbbox((x, current_y), line, font=font)
        current_y = bbox[3] + line_gap
      current_y += line_gap
    return current_y


def create_slide_image(theme: CarouselTheme, slide_index: int, out_path: Path) -> None:
    image = Image.new("RGB", (SLIDE_WIDTH, SLIDE_HEIGHT), color=theme.background)
    draw = ImageDraw.Draw(image)
    accent_rgb = to_rgb(theme.accent)
    text_color = theme.text
    muted_color = theme.muted

    draw.rounded_rectangle((70, 60, 1010, 1290), radius=44, outline=accent_rgb, width=3)
    draw.rectangle((70, 60, 180, 1290), fill=accent_rgb)
    draw.ellipse((860, 118, 970, 228), outline=accent_rgb, width=6)
    draw.line((220, 240, 880, 240), fill=accent_rgb, width=3)

    if slide_index == 0:
        draw.text((220, 126), theme.kicker.upper(), font=SMALL_FONT, fill=theme.accent)
        y = add_text_block(
            draw,
            text=theme.title,
            xy=(220, 300),
            font=TITLE_FONT,
            fill=text_color,
            max_width=620,
            line_gap=16,
        )
        add_text_block(
            draw,
            text=theme.subtitle,
            xy=(220, y + 30),
            font=BODY_FONT,
            fill=muted_color,
            max_width=610,
            line_gap=12,
        )
        draw.rounded_rectangle((220, 980, 820, 1090), radius=28, fill=theme.accent)
        draw.text((260, 1016), "Edit headline, subtitle, and brand line", font=BODY_FONT, fill=theme.background)
    elif slide_index == 1:
        draw.text((220, 126), theme.kicker.upper(), font=SMALL_FONT, fill=theme.accent)
        draw.text((220, 285), theme.insight_title, font=TITLE_FONT, fill=text_color)
        current_y = 450
        for bullet in theme.bullets:
            draw.ellipse((222, current_y + 12, 246, current_y + 36), fill=theme.accent)
            add_text_block(
                draw,
                text=bullet,
                xy=(270, current_y),
                font=BODY_FONT,
                fill=text_color,
                max_width=560,
                line_gap=10,
            )
            current_y += 150
        draw.rounded_rectangle((220, 1090, 860, 1210), radius=26, outline=accent_rgb, width=4)
        draw.text((260, 1132), "Swap bullets for your own framework or key points", font=BODY_FONT, fill=muted_color)
    else:
        draw.text((220, 126), theme.kicker.upper(), font=SMALL_FONT, fill=theme.accent)
        draw.text((220, 285), theme.cta_title, font=TITLE_FONT, fill=text_color)
        current_y = 470
        for index, line in enumerate(theme.cta_lines, start=1):
            draw.rounded_rectangle((220, current_y, 870, current_y + 120), radius=22, fill="#FFFFFF")
            draw.text((255, current_y + 38), f"{index:02d}", font=BODY_FONT, fill=theme.accent)
            draw.text((350, current_y + 36), line, font=BODY_FONT, fill=text_color)
            current_y += 150
        draw.text((220, 1140), "Duplicate slides, replace copy, and publish as a carousel.", font=BODY_FONT, fill=muted_color)

    image.save(out_path)


def add_text_box(slide, left, top, width, height, text, font_size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*to_rgb(color))
    return box


def create_slide_presentation(theme: CarouselTheme, slide, slide_index: int) -> None:
    fill_rgb = RGBColor(*to_rgb(theme.background))
    accent_rgb = RGBColor(*to_rgb(theme.accent))

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = fill_rgb

    left_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(0.55), Inches(0.9), Inches(11.4))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = accent_rgb
    left_bar.line.color.rgb = accent_rgb

    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(0.55), Inches(8.7), Inches(11.4))
    frame.fill.background()
    frame.line.color.rgb = accent_rgb
    frame.line.width = Pt(2)

    add_text_box(slide, Inches(1.95), Inches(1.1), Inches(5.8), Inches(0.45), theme.kicker.upper(), 18, theme.accent, bold=True)

    if slide_index == 0:
        add_text_box(slide, Inches(1.95), Inches(2.2), Inches(5.9), Inches(2.8), theme.title, 28, theme.text, bold=True)
        add_text_box(slide, Inches(1.95), Inches(5.2), Inches(5.8), Inches(2.0), theme.subtitle, 18, theme.muted)
        pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.95), Inches(8.5), Inches(5.35), Inches(0.9))
        pill.fill.solid()
        pill.fill.fore_color.rgb = accent_rgb
        pill.line.color.rgb = accent_rgb
        add_text_box(slide, Inches(2.25), Inches(8.73), Inches(4.8), Inches(0.4), "Edit headline, subtitle, and brand line", 16, theme.background, bold=True)
    elif slide_index == 1:
        add_text_box(slide, Inches(1.95), Inches(2.2), Inches(5.8), Inches(0.8), theme.insight_title, 28, theme.text, bold=True)
        top = 3.7
        for bullet in theme.bullets:
            dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(2.0), Inches(top + 0.1), Inches(0.18), Inches(0.18))
            dot.fill.solid()
            dot.fill.fore_color.rgb = accent_rgb
            dot.line.color.rgb = accent_rgb
            add_text_box(slide, Inches(2.35), Inches(top), Inches(5.1), Inches(0.7), bullet, 20, theme.text)
            top += 1.3
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.95), Inches(8.95), Inches(5.7), Inches(1.0))
        card.fill.background()
        card.line.color.rgb = accent_rgb
        card.line.width = Pt(2)
        add_text_box(slide, Inches(2.2), Inches(9.23), Inches(5.15), Inches(0.5), "Swap bullets for your own framework or key points", 16, theme.muted)
    else:
        add_text_box(slide, Inches(1.95), Inches(2.2), Inches(5.8), Inches(0.8), theme.cta_title, 28, theme.text, bold=True)
        top = 3.8
        for index, line in enumerate(theme.cta_lines, start=1):
            card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.95), Inches(top), Inches(5.85), Inches(0.95))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = accent_rgb
            add_text_box(slide, Inches(2.2), Inches(top + 0.26), Inches(0.8), Inches(0.4), f"{index:02d}", 20, theme.accent, bold=True)
            add_text_box(slide, Inches(3.05), Inches(top + 0.23), Inches(4.4), Inches(0.4), line, 18, theme.text)
            top += 1.2
        add_text_box(
            slide,
            Inches(1.95),
            Inches(9.1),
            Inches(5.8),
            Inches(1.0),
            "Duplicate slides, replace copy, and publish as a carousel.",
            16,
            theme.muted,
        )


def build_pptx(themes: list[CarouselTheme], out_path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(INSTA_RATIO[0])
    presentation.slide_height = Inches(INSTA_RATIO[1])
    blank_layout = presentation.slide_layouts[6]

    for theme in themes:
        for slide_index in range(3):
            slide = presentation.slides.add_slide(blank_layout)
            create_slide_presentation(theme, slide, slide_index)

    presentation.save(out_path)


def build_quickstart_pdf(out_path: Path, pack_title: str) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    styles = getSampleStyleSheet()
    styles.add(
        ParagraphStyle(
            name="BodySmall",
            parent=styles["BodyText"],
            fontName="Helvetica",
            fontSize=10,
            leading=14,
            textColor=colors.HexColor("#3D3A36"),
        )
    )
    story = [
        Paragraph(pack_title, styles["Title"]),
        Spacer(1, 0.15 * inch),
        Paragraph("Quick Start Guide", styles["Heading2"]),
        Spacer(1, 0.1 * inch),
        Paragraph(
            "Open the included PowerPoint file, duplicate any slide you want to reuse, and replace the placeholder copy with your own message.",
            styles["BodySmall"],
        ),
        Spacer(1, 0.12 * inch),
        Paragraph("Recommended workflow", styles["Heading3"]),
        Paragraph(
            "1. Pick a 3-slide sequence that matches the idea you want to teach. "
            "2. Update the headline, proof points, and CTA. "
            "3. Export the edited slides as PNGs at 1080 x 1350 for Instagram.",
            styles["BodySmall"],
        ),
        Spacer(1, 0.12 * inch),
        Paragraph("Usage notes", styles["Heading3"]),
        Paragraph(
            "You may use these templates for your own brand and for client delivery work. "
            "Do not resell, redistribute, or repackage the original template files as a competing template product.",
            styles["BodySmall"],
        ),
        Spacer(1, 0.12 * inch),
        Paragraph("Included files", styles["Heading3"]),
        Paragraph(
            "- Editable PowerPoint deck with 30 slides<br/>"
            "- PDF quick-start guide<br/>"
            "- Store preview image<br/>"
            "- Simple license note",
            styles["BodySmall"],
        ),
    ]
    document = SimpleDocTemplate(str(out_path), pagesize=letter, leftMargin=0.75 * inch, rightMargin=0.75 * inch)
    document.build(story)


def render_pdf_preview(pdf_path: Path, output_dir: Path) -> list[Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    document = fitz.open(pdf_path)
    rendered: list[Path] = []
    for page_number in range(document.page_count):
        page = document.load_page(page_number)
        pixmap = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
        path = output_dir / f"quick-start-{page_number + 1:02d}.png"
        pixmap.save(path)
        rendered.append(path)
    document.close()
    return rendered


def create_cover(preview_paths: list[Path], out_path: Path, title: str, subtitle: str, *, square: bool = False) -> None:
    width, height = (1200, 1200) if square else (1536, 1024)
    canvas = Image.new("RGB", (width, height), "#F5F1EB")
    draw = ImageDraw.Draw(canvas)
    shadow = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    shadow_draw = ImageDraw.Draw(shadow)

    sample_paths = preview_paths[:4]
    card_width = 250 if square else 260
    card_height = 312 if square else 324
    gap = 30
    total_width = len(sample_paths) * card_width + (len(sample_paths) - 1) * gap
    start_x = (width - total_width) // 2
    top_y = 170 if square else 180

    for index, path in enumerate(sample_paths):
        card = Image.open(path).convert("RGB").resize((card_width, card_height))
        x = start_x + index * (card_width + gap)
        y = top_y
        shadow_draw.rounded_rectangle((x + 10, y + 18, x + card_width + 10, y + card_height + 18), radius=22, fill=(0, 0, 0, 55))
        canvas.paste(card, (x, y))
        draw.rounded_rectangle((x, y, x + card_width, y + card_height), radius=22, outline="#817161", width=4)

    canvas = Image.alpha_composite(canvas.convert("RGBA"), shadow.filter(ImageFilter.GaussianBlur(radius=12))).convert("RGB")
    draw = ImageDraw.Draw(canvas)
    title_font = find_font(
        [
            "C:/Windows/Fonts/georgiab.ttf",
            "/usr/share/fonts/truetype/dejavu/DejaVuSerif-Bold.ttf",
        ],
        46 if square else 54,
    )
    subtitle_font = find_font(
        [
            "C:/Windows/Fonts/arial.ttf",
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        ],
        26 if square else 30,
    )
    title_y = 620 if square else 610
    draw.text((120, title_y), title, font=title_font, fill="#1E1C1A")
    subtitle_lines = textwrap.wrap(subtitle, width=48 if square else 60)
    current_y = title_y + 80
    for line in subtitle_lines:
        draw.text((120, current_y), line, font=subtitle_font, fill="#67615B")
        current_y += 40
    badge_text = "Editable PowerPoint templates"
    draw.rounded_rectangle((120, height - 150, 640, height - 80), radius=26, fill="#7F6A58")
    draw.text((160, height - 132), badge_text, font=subtitle_font, fill="#F6F0EA")
    canvas.save(out_path)


def write_license(out_path: Path) -> None:
    out_path.write_text(
        "\n".join(
            [
                "ImonEngine template license",
                "",
                "You may use these templates for your own brand and for client work.",
                "You may edit the layouts, copy, colors, and imagery for finished marketing assets.",
                "Do not resell, redistribute, or repackage the original template files as a competing template product.",
            ]
        ),
        encoding="utf-8",
    )


def update_pack_state(pack_dir: Path, pack_data: dict) -> None:
    runtime_dir = pack_dir.parent.parent
    state_path = runtime_dir / "state" / "assetPacks.json"
    if not state_path.exists():
        return
    packs = json.loads(state_path.read_text(encoding="utf-8"))
    for index, candidate in enumerate(packs):
        if candidate.get("id") == pack_data["id"]:
            packs[index] = pack_data
            break
    state_path.write_text(json.dumps(packs, indent=2) + "\n", encoding="utf-8")


def write_pack_artifacts(pack_dir: Path, pack_data: dict) -> None:
    listing_lines = [
        f"# {pack_data['title']}",
        "",
        f"Marketplace: {pack_data['marketplace']}",
        f"Status: {pack_data['status']}",
        f"Suggested price: ${pack_data['suggestedPrice']}",
        f"Price test points: {', '.join(f'${value}' for value in pack_data['priceVariants'])}",
        "",
        "## Summary",
        pack_data["shortDescription"],
        "",
        "## Description",
        pack_data["description"],
        "",
        "## Deliverables",
        *[f"- {item}" for item in pack_data["deliverables"]],
        "",
        "## Tags",
        ", ".join(f"`{tag}`" for tag in pack_data["tags"]),
    ]
    (pack_dir / "listing.md").write_text("\n".join(listing_lines) + "\n", encoding="utf-8")
    (pack_dir / "manifest.json").write_text(json.dumps(pack_data, indent=2) + "\n", encoding="utf-8")
    gumroad_draft = [
        f"# {pack_data['title']}",
        "",
        f"Suggested price: ${pack_data['suggestedPrice']}",
        f"Price tests: {', '.join(f'${value}' for value in pack_data['priceVariants'])}",
        "",
        "## Short Description",
        pack_data["shortDescription"],
        "",
        "## Full Description",
        pack_data["description"],
        "",
        "## Deliverables",
        *[f"- {item}" for item in pack_data["deliverables"]],
        "",
        "## Tags",
        ", ".join(pack_data["tags"]),
    ]
    (pack_dir / "gumroad" / "product-draft.md").write_text("\n".join(gumroad_draft) + "\n", encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--pack-dir", required=True)
    args = parser.parse_args()

    pack_dir = Path(args.pack_dir).resolve()
    manifest_path = pack_dir / "manifest.json"
    pack = json.loads(manifest_path.read_text(encoding="utf-8"))

    previews_dir = pack_dir / "assets" / "final" / "preview-images"
    covers_dir = pack_dir / "covers"
    gumroad_dir = pack_dir / "gumroad"
    product_files_dir = gumroad_dir / "product-files"
    tmp_pdf_dir = pack_dir / "tmp" / "pdfs"

    for directory in [previews_dir, covers_dir, gumroad_dir, product_files_dir, tmp_pdf_dir]:
        directory.mkdir(parents=True, exist_ok=True)

    preview_paths: list[Path] = []
    for theme in THEMES:
        for slide_index in range(3):
            slide_no = len(preview_paths) + 1
            path = previews_dir / f"neutral-carousel-template-{slide_no:02d}.png"
            create_slide_image(theme, slide_index, path)
            preview_paths.append(path)

    pptx_path = product_files_dir / "neutral-instagram-carousel-template-pack.pptx"
    build_pptx(THEMES, pptx_path)

    quickstart_pdf = product_files_dir / "quick-start-guide.pdf"
    build_quickstart_pdf(quickstart_pdf, pack["title"])
    render_pdf_preview(quickstart_pdf, tmp_pdf_dir)

    overview_path = product_files_dir / "template-overview.png"
    create_cover(
        preview_paths,
        overview_path,
        pack["title"],
        "Neutral editorial carousel templates for creators, coaches, and solo businesses.",
        square=False,
    )
    cover_one = covers_dir / "cover-01.png"
    cover_two = covers_dir / "cover-02.png"
    thumb = covers_dir / "thumbnail-square.png"
    create_cover(preview_paths, cover_one, pack["title"], "30 editable slides in a soft editorial style.", square=False)
    create_cover(preview_paths[4:] + preview_paths[:4], cover_two, "Instagram Carousel Templates", "Built for calm, premium-looking creator marketing.", square=False)
    create_cover(preview_paths, thumb, "Neutral Carousel Templates", "Editable PowerPoint layouts for creators.", square=True)

    license_path = product_files_dir / "LICENSE.txt"
    write_license(license_path)

    zip_path = gumroad_dir / "neutral-instagram-carousel-template-pack.zip"
    with zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for path in [pptx_path, quickstart_pdf, overview_path, license_path]:
            archive.write(path, arcname=path.name)

    pack.update(
        {
            "title": "Neutral Instagram Carousel Template Pack",
            "shortDescription": "A neutral carousel template pack for creators who want clean posts without custom design work.",
            "description": "A ready-to-use Instagram carousel template pack for creators, coaches, and solo businesses. The layouts use a restrained neutral palette so the pack feels premium instead of generic. The included PowerPoint deck makes each slide editable for faster posting and client delivery.",
            "suggestedPrice": 12,
            "priceVariants": [9, 12, 15],
            "tags": [
                "instagram carousel template",
                "editable powerpoint template",
                "neutral social template",
                "creator marketing",
                "gumroad digital download",
            ],
            "deliverables": [
                "30 editable Instagram carousel templates in PowerPoint format",
                "Quick-start guide PDF",
                "2 Gumroad cover images and 1 square thumbnail",
                "Simple commercial-use license note",
            ],
            "status": "ready_for_upload",
            "updatedAt": now_iso(),
        }
    )

    write_pack_artifacts(pack_dir, pack)
    update_pack_state(pack_dir, pack)
    print(
        json.dumps(
            {
                "packId": pack["id"],
                "status": pack["status"],
                "zipPath": str(zip_path),
                "pptxPath": str(pptx_path),
                "previewCount": len(preview_paths),
            },
            indent=2,
        )
    )


if __name__ == "__main__":
    main()
